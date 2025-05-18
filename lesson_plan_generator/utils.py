"""Utility functions for lesson plan generation."""
import os
import logging
from docx import Document
from pptx import Presentation
from django.conf import settings

logger = logging.getLogger(__name__)

def ensure_dir(file_path):
    """确保目录存在，如果不存在则创建"""
    directory = os.path.dirname(file_path)
    if not os.path.exists(directory):
        os.makedirs(directory)
        logger.info(f"Created directory: {directory}")

def get_file_path(lesson_id, extension):
    """获取文件的完整路径"""
    return os.path.join(settings.MEDIA_ROOT, 'lessons', f"{lesson_id}.{extension}")

def generate_attachments(lesson, ai_content=None):
    """
    Generate Word and PowerPoint attachments for a lesson plan.
    
    Args:
        lesson: LessonPlan instance to generate attachments for
        ai_content: Optional AI generated content dictionary
    """
    logger.info(f"Starting to generate attachments for lesson {lesson.id}")
    
    # 生成Word文档
    try:
        doc = Document()
        doc.add_heading(f"{lesson.get_subject_display()}教案", 0)
        
        # 添加基本信息
        doc.add_heading('基本信息', level=1)
        doc.add_paragraph(f'年级：{lesson.get_grade_display()}')
        doc.add_paragraph(f'课时：{lesson.duration}')
        doc.add_paragraph(f'教学方式：{lesson.get_teaching_style_display()}')
        
        if ai_content:
            # 使用AI生成的内容
            sections = ai_content['sections']
            
            # 添加教学目标
            doc.add_heading('教学目标', level=1)
            doc.add_paragraph(sections['objectives'])
            
            # 添加教学步骤
            doc.add_heading('教学步骤', level=1)
            doc.add_paragraph(sections['steps'])
            
            # 添加重难点讲解策略
            doc.add_heading('重难点讲解策略', level=1)
            doc.add_paragraph(sections['strategies'])
            
            # 添加课堂互动设计
            doc.add_heading('课堂互动设计', level=1)
            doc.add_paragraph(sections['interactions'])
            
            # 添加课堂提问和检测题目
            doc.add_heading('课堂提问和检测题目', level=1)
            doc.add_paragraph(sections['questions'])
            
            # 添加课后作业建议
            doc.add_heading('课后作业建议', level=1)
            doc.add_paragraph(sections['homework'])
        else:
            # 使用原始内容
            doc.add_heading('教学目标', level=1)
            doc.add_paragraph(lesson.objectives)
            
            doc.add_heading('教学重点', level=1)
            for point in lesson.key_points.split(','):
                if point.strip():
                    doc.add_paragraph(point.strip(), style='List Bullet')
            
            if lesson.difficulties:
                doc.add_heading('教学难点', level=1)
                doc.add_paragraph(lesson.difficulties)
            
            if lesson.student_profile:
                doc.add_heading('学生情况', level=1)
                doc.add_paragraph(lesson.student_profile)

        # 确保目录存在并保存文件
        file_path = get_file_path(lesson.id, 'docx')
        ensure_dir(file_path)
        logger.info(f"Saving Word document to {file_path}")
        doc.save(file_path)
        logger.info("Word document saved successfully")
    except Exception as e:
        logger.error(f"生成Word文档失败：{str(e)}", exc_info=True)
        raise

    # 生成PPT
    try:
        prs = Presentation()
        
        # 标题页
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = f"{lesson.get_subject_display()}教案"
        subtitle.text = f"{lesson.get_grade_display()} - {lesson.duration}"
        
        if ai_content:
            # 使用AI生成的PPT大纲
            sections = ai_content['sections']
            
            # 教学目标页
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            body = slide.placeholders[1]
            title.text = "教学目标"
            body.text = sections['objectives']
            
            # 教学步骤页
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            body = slide.placeholders[1]
            title.text = "教学步骤"
            body.text = sections['steps']
            
            # 重难点讲解策略页
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            body = slide.placeholders[1]
            title.text = "重难点讲解策略"
            body.text = sections['strategies']
            
            # 课堂互动设计页
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            body = slide.placeholders[1]
            title.text = "课堂互动设计"
            body.text = sections['interactions']
            
            # PPT大纲页
            if sections.get('ppt_outline'):
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                title = slide.shapes.title
                body = slide.placeholders[1]
                title.text = "PPT大纲"
                body.text = sections['ppt_outline']
        else:
            # 使用原始内容
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            body = slide.placeholders[1]
            title.text = "教学目标"
            body.text = lesson.objectives
            
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            body = slide.placeholders[1]
            title.text = "教学重点"
            points = lesson.key_points.split(',')
            body.text = '\n'.join(f"• {point.strip()}" for point in points if point.strip())

        # 确保目录存在并保存文件
        file_path = get_file_path(lesson.id, 'pptx')
        ensure_dir(file_path)
        logger.info(f"Saving PowerPoint to {file_path}")
        prs.save(file_path)
        logger.info("PowerPoint saved successfully")
    except Exception as e:
        logger.error(f"生成PPT失败：{str(e)}", exc_info=True)
        raise 